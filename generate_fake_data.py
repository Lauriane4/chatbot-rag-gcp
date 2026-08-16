import os
import pandas as pd
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors

BASE_DIR = "data_gustocraft"

def create_pdf(path, title, ref_doc, version, date_doc, author, approver, content_list):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    doc = SimpleDocTemplate(path, pagesize=letter, rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36)
    styles = getSampleStyleSheet()

    header_table_data = [
        [
            Paragraph("<b>GUSTOO</b>", ParagraphStyle('H1', fontName='Helvetica-Bold', fontSize=14, leading=16, textColor=colors.HexColor("#1A365D"))),
            Paragraph(f"<b>Réf. document :</b> {ref_doc}<br/><b>Version :</b> {version}<br/><b>Date :</b> {date_doc}", ParagraphStyle('H2', fontName='Helvetica', fontSize=8, leading=10))
        ],
        [
            Paragraph(f"<b>{title}</b>", ParagraphStyle('H3', fontName='Helvetica-Bold', fontSize=11, leading=13, textColor=colors.HexColor("#2B6CB0"))),
            Paragraph(f"<b>Rédigé par :</b> {author}<br/><b>Approuvé par :</b> {approver}", ParagraphStyle('H4', fontName='Helvetica', fontSize=8, leading=10))
        ],
        [
            Paragraph("<b>CONFIDENTIEL — DIFFUSION RESTREINTE</b>", ParagraphStyle('H5', fontName='Helvetica-Bold', fontSize=8, leading=9, textColor=colors.HexColor("#C53030"))),
            ""
        ]
    ]
    t_header = Table(header_table_data, colWidths=[320, 220])
    t_header.setStyle(TableStyle([
        ('BOX', (0,0), (-1,-1), 1, colors.HexColor("#CBD5E0")),
        ('INNERGRID', (0,0), (-1,-1), 0.5, colors.HexColor("#E2E8F0")),
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#F7FAFC")),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('PADDING', (0,0), (-1,-1), 4),
        ('SPAN', (0,2), (1,2)),
        ('ALIGN', (0,2), (1,2), 'CENTER'),
        ('BACKGROUND', (0,2), (1,2), colors.HexColor("#FFF5F5")),
    ]))

    h2_style = ParagraphStyle('DocH2', parent=styles['Heading2'], fontSize=11, leading=14, textColor=colors.HexColor("#1A365D"), spaceBefore=8, spaceAfter=4)
    body_style = ParagraphStyle('DocBody', parent=styles['Normal'], fontSize=8.5, leading=11.5, textColor=colors.HexColor("#2D3748"), spaceAfter=5)

    elements = [t_header, Spacer(1, 10)]

    for item_type, item_content in content_list:
        if item_type == "h2":
            elements.append(Paragraph(f"<b>{item_content}</b>", h2_style))
        elif item_type == "p":
            elements.append(Paragraph(item_content, body_style))
        elif item_type == "table":
            t_data = []
            for row in item_content:
                row_data = [Paragraph(str(cell), ParagraphStyle('TCell', fontSize=7.5, leading=9.5)) for cell in row]
                t_data.append(row_data)
            t = Table(t_data)
            t.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.HexColor("#EDF2F7")),
                ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
                ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor("#CBD5E0")),
                ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
                ('PADDING', (0,0), (-1,-1), 3.5),
            ]))
            elements.append(t)
            elements.append(Spacer(1, 6))

    doc.build(elements)

def create_docx(path, title, ref_doc, version, date_doc, author, approver, sections):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    doc = Document()
    
    # En-tête officiel
    header_table = doc.add_table(rows=2, cols=2)
    header_table.style = 'Table Grid'
    header_table.cell(0, 0).text = f"GUSTOO\n{title}"
    header_table.cell(0, 1).text = f"Réf. document : {ref_doc}\nVersion : {version}\nDate : {date_doc}"
    header_table.cell(1, 0).text = "CONFIDENTIEL — DIFFUSION RESTREINTE"
    header_table.cell(1, 1).text = f"Rédigé par : {author}\nApprouvé par : {approver}"
    doc.add_paragraph()

    for sec_title, text_blocks in sections:
        if sec_title:
            doc.add_heading(sec_title, level=2)
        for block in text_blocks:
            if isinstance(block, list):
                table = doc.add_table(rows=len(block), cols=len(block[0]))
                table.style = 'Light Shading Accent 1'
                for r_idx, row in enumerate(block):
                    for c_idx, val in enumerate(row):
                        table.cell(r_idx, c_idx).text = str(val)
                doc.add_paragraph()
            else:
                doc.add_paragraph(block)
    doc.save(path)

def create_pptx(path, slides_data):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs = Presentation()
    for title_text, bullets in slides_data:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_text
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        for idx, bullet in enumerate(bullets):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.text = bullet
            p.font.size = Pt(14)
    prs.save(path)

def create_excel(path, sheets_dict):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with pd.ExcelWriter(path, engine='openpyxl') as writer:
        for sheet_name, df in sheets_dict.items():
            df.to_excel(writer, sheet_name=sheet_name, index=False)

def generer_documents():
    print("🚀 Génération complète des documents GustoCraft / GUSTOO...")

    # ==========================================
    # 00_Direction_et_Comex
    # ==========================================
    create_docx(
        f"{BASE_DIR}/00_Direction_et_Comex/CR_COMEX_Janvier_2026.docx",
        "Compte-Rendu du Comité Exécutif - Séance du 15 Janvier 2026",
        "CR-COMEX-2026-01", "1.0", "15 janvier 2026",
        "Nadia BELKACEM — Assistante de Direction",
        "Camille ROUSSEL — Directrice Générale",
        [
            ("1. Synthèse Opérationnelle & Performance 2025", [
                "Présents : Camille ROUSSEL (DG), Julien MERCIER (DGA), Nathalie PERRIN (DAF), Catherine LEFEVRE (DRH), Marc-Antoine PELLETIER (Dir. Industriel), Sophie LAMBERT (Dir. R&D), Hélène GIRAUD (Dir. QSE), Philippe MOREL (Dir. Achats), Isabelle ROBIN (Dir. Commerciale).",
                "Le chiffre d'affaires consolidé 2025 s'établit à 48,2 M€ (+6,5% vs 2024), porté par le Client Alpha Distribution (18,5 M€) et le Client Bêta Foodservice (12,4 M€).",
                "Marge brute globale : 34,8% (recul de 1,2 pt lié aux cours du beurre et des huiles végétales)."
            ]),
            ("2. Décisions Stratégiques & Arbitrages Industriels", [
                "Validation du CAPEX de 1,45 M€ pour l'automatisation de l'operculage sur la Ligne 3 du site d'Amiens (Usine Nord).",
                "Second sourcing obligatoire : Claire VASSEUR doit finaliser la qualification du Fournisseur B en secours du Fournisseur A pour la farine T45 avant le 31 mars 2026 selon la procédure PROC_ACH_007.",
                "Plan RSE Emballage : Validation de la cible 2026 fixée par Thomas WEBER de -25% de plastique vierge (barquettes rPET 450µm)."
            ])
        ]
    )

    create_pptx(
        f"{BASE_DIR}/00_Direction_et_Comex/PRESENTATION_Plan_Strategique_2026-2028.pptx",
        [
            ("GUSTOO — Plan Stratégique 2026-2028", [
                "Comité de Direction — Février 2026",
                "Présenté par Camille ROUSSEL (DG) et Julien MERCIER (COO)",
                "CONFIDENTIEL — DIFFUSION RESTREINTE"
            ]),
            ("Axe 1 : Excellence Industrielle & TRS Usines", [
                "Objectif TRS Usine Nord (Amiens) et Usine Sud (Avignon) : ≥ 82% (AIDE_CALCUL_05).",
                "Automatisation complète des lignes d'operculage et de conditionnement sous vide.",
                "Déploiement du SMED sur les changements de format brioches et sauces."
            ]),
            ("Axe 2 : R&D & Clean Label", [
                "Suppression totale des émulsifiants de synthèse d'ici fin 2026 (Karim OUALI & Elodie FONTAINE).",
                "Intégration des substituts naturels (Allulose, protéine de féverole du Fournisseur C).",
                "Renouvellement des certifications IFS Food v8 avec un score visé supérieur à 96%."
            ])
        ]
    )

    # ==========================================
    # 01_Procedures_et_Guidelines (Compléments)
    # ==========================================
    create_docx(
        f"{BASE_DIR}/01_Procedures_et_Guidelines/GUIDELINES_DEV_004_Creation_et_Validation_Nouvel_Emballage.docx",
        "Création et Validation d'un Nouvel Emballage",
        "GUIDELINES-DEV-004", "1.2", "18 janvier 2026",
        "Léa SIMON — Ingénieure Packaging",
        "Thomas WEBER — Resp. Packaging & DevPack",
        [
            ("1. Objet", [
                "Ces lignes directrices décrivent le processus de création et de validation d'un nouvel emballage pour les produits GUSTOO, depuis l'expression du besoin jusqu'à la validation industrielle définitive sur les sites d'Amiens et d'Avignon."
            ]),
            ("2. Acteurs impliqués", [
                [
                    ["Rôle", "Acteur", "Contribution au projet"],
                    ["Resp. Packaging & DevPack", "Thomas WEBER", "Pilote le projet de bout en bout et arbitre les choix techniques."],
                    ["Ingénieure Packaging", "Léa SIMON", "Réalise les tests techniques (étanchéité, perméabilité) et le CDC matériau."],
                    ["Acheteur Emballages", "David ROUX", "Sourcing fournisseur, négociation et qualification (PROC_ACH_007)."],
                    ["Resp. Assurance Qualité", "Nicolas BERTIN", "Valide la conformité réglementaire (CE 1935/2004) et alimentarité."],
                    ["Directeur Industriel", "Marc-Antoine PELLETIER", "Valide la faisabilité industrielle sur les lignes existantes."]
                ]
            ]),
            ("3. Délais indicatifs par étape", [
                [
                    ["Étape", "Délai indicatif", "Livrable obligatoire"],
                    ["Cahier des charges fonctionnel", "2 à 3 semaines", "Fiche FICHE-EMB-V3 validée en comité DevPack"],
                    ["Sourcing et qualification fournisseur", "4 à 8 semaines", "Dossier PROC_ACH_007 complet avec David ROUX"],
                    ["Tests techniques laboratoire", "3 à 4 semaines", "Rapport TEST_Etancheite_et_Permeabilite_O2"],
                    ["Test de vieillissement produit", "Jusqu'à 90 jours", "Analyses microbiologiques J0, J+30, J+60, J+90"],
                    ["Essai pilote sur ligne industrielle", "1 à 2 semaines", "Protocole GUIDELINES_RD_019 validé"]
                ]
            ]),
            ("4. Historique des révisions", [
                [
                    ["Version", "Date", "Nature de la modification", "Auteur"],
                    ["1.0", "12/02/2023", "Création initiale", "T. WEBER"],
                    ["1.1", "30/08/2024", "Ajout du volet test de vieillissement accéléré", "L. SIMON"],
                    ["1.2", "18/01/2026", "Précision des délais indicatifs et renvoi GUIDE_CALCUL_03", "T. WEBER"]
                ]
            ])
        ]
    )

    create_docx(
        f"{BASE_DIR}/01_Procedures_et_Guidelines/GUIDELINES_RD_019_Demande_Essai_Pilote_Ligne_Industrielle.docx",
        "Demande d'Essai Pilote sur Ligne Industrielle",
        "GUIDELINES-RD-019", "1.2", "20 janvier 2026",
        "Sophie LAMBERT — Directrice R&D & Innovation",
        "Marc-Antoine PELLETIER — Directeur Industriel",
        [
            ("1. Objet", [
                "Ces lignes directrices encadrent la demande, la planification et le déroulement d'un essai pilote sur ligne industrielle, réalisé par la R&D en amont d'un lancement de nouvelle formulation ou d'un nouvel emballage."
            ]),
            ("2. Acteurs et responsabilités", [
                [
                    ["Rôle", "Responsable", "Mission dans l'essai pilote"],
                    ["Ingénieur R&D demandeur", "Yann PICARD / Manon DUBUISSON", "Formalise la demande et définit le protocole expérimental."],
                    ["Resp. Formulation", "Karim OUALI / Elodie FONTAINE", "Valide le protocole technique et supervise l'essai."],
                    ["Directrice Usine Nord", "Odile FABRE", "Valide le créneau sur Amiens (Ligne 1 ou Ligne 2)."],
                    ["Directeur Usine Sud", "Grégory ANDRE", "Valide le créneau sur Avignon (Ligne 3 Sauces)."],
                    ["Resp. Maintenance Usine Nord", "Vincent LACROIX", "Assure le calage machine et l'outillage spécifique."]
                ]
            ]),
            ("3. Créneaux réservés et délais de soumission", [
                "Toute demande doit être déposée via le formulaire DEMANDE-ESSAI-PILOTE au moins 3 semaines avant la date souhaitée.",
                "Créneau dédié Usine Nord (Amiens) : Premier vendredi après-midi de chaque mois.",
                "Créneau dédié Usine Sud (Avignon) : Troisième mercredi matin de chaque mois.",
                "Le temps de ligne mobilisé est imputé au budget R&D selon le tarif horaire interne communiqué par Nathalie PERRIN (DAF)."
            ]),
            ("4. Historique des révisions", [
                [
                    ["Version", "Date", "Nature de la modification", "Auteur"],
                    ["1.0", "10/05/2022", "Création initiale", "S. LAMBERT"],
                    ["1.1", "14/11/2023", "Ajout des créneaux mensuels dédiés par site", "K. OUALI"],
                    ["1.2", "20/01/2026", "Ajout du volet imputation des coûts au budget R&D", "S. LAMBERT"]
                ]
            ])
        ]
    )

    # ==========================================
    # 02_Guides_Calculs_et_Matrices
    # ==========================================
    create_pdf(
        f"{BASE_DIR}/02_Guides_Calculs_et_Matrices/GUIDE_CALCUL_01_Pourcentages_Boulanger_et_Hydratation.pdf",
        "Calcul des Pourcentages du Boulanger et Taux d'Hydratation",
        "GUIDE-CALCUL-01", "1.1", "12 janvier 2026",
        "Karim OUALI — Resp. Formulation Boulangerie",
        "Sophie LAMBERT — Directrice R&D & Innovation",
        [
            ("h2", "1. Principe du Pourcentage du Boulanger (Baker's Percentage)"),
            ("p", "En panification et viennoiserie industrielle chez GUSTOO, la masse totale de farine représente TOUJOURS la base 100 %. Tous les ingrédients sont dosés en proportion de cette masse de référence."),
            ("p", "Formule : % Ingrédient = [ Masse Ingrédient (kg) / Masse Totale Farine (kg) ] × 100"),
            ("h2", "2. Règle de Calcul du Taux d'Hydratation Réel"),
            ("p", "L'eau totale inclut l'eau osmosée, 90% du poids du lait liquide et 75% du poids des œufs entiers pasteurisés (MP-OEU-002 du Fournisseur C)."),
            ("p", "Formule : Eau Totale (kg) = Eau_Coulage + (0.90 × Lait) + (0.75 × Oeufs)"),
            ("p", "Taux d'Hydratation (%) = [ Eau Totale (kg) / Masse Totale Farine (kg) ] × 100"),
            ("h2", "3. Formulation Standard Brioche Feuilletée (Pétrin Pilote 100 kg Farine)"),
            ("table", [
                ["Ingrédient", "Code MP", "Fournisseur Agréé", "% Boulanger", "Poids (kg)", "Apport Eau"],
                ["Farine Blé T45", "MP-FAR-001", "Fournisseur A (Fournisseur B)", "100.0 %", "100.00 kg", "0.00 kg"],
                ["Beurre Tourage 84%", "MP-BEU-004", "Fournisseur D Laiterie", "35.0 %", "35.00 kg", "5.60 kg"],
                ["Eau filtrée osmosée", "MP-EAU-000", "Réseau Interne Amiens", "28.0 %", "28.00 kg", "28.00 kg"],
                ["Oeufs entiers past.", "MP-OEU-002", "Fournisseur C Ingrédients", "30.0 %", "30.00 kg", "22.50 kg"],
                ["Sucre semoule fin", "MP-SUC-001", "Fournisseur E Sucres", "15.0 %", "15.00 kg", "0.00 kg"],
                ["Levure liquide", "MP-LEV-001", "Fournisseur F Levures", "4.5 %", "4.50 kg", "3.15 kg"],
                ["Sel fin de mer", "MP-SEL-001", "Fournisseur G Sel", "2.0 %", "2.00 kg", "0.00 kg"],
                ["Total Pétrin", "—", "—", "214.5 %", "214.50 kg", "59.25 kg (TH: 59.25%)"]
            ]),
            ("h2", "4. Historique des révisions"),
            ("table", [
                ["Version", "Date", "Nature de la modification", "Auteur"],
                ["1.0", "14/04/2023", "Création initiale", "K. OUALI"],
                ["1.1", "12/01/2026", "Précision de l'apport en eau des ovoproduits", "Y. PICARD"]
            ])
        ]
    )

    create_pdf(
        f"{BASE_DIR}/02_Guides_Calculs_et_Matrices/GUIDE_CALCUL_04_Bareme_Frais_Kilometriques_et_Repas_2026.pdf",
        "Barème des Frais Professionnels, Déplacements et Télétravail",
        "GUIDE-CALCUL-04", "2.0", "2 janvier 2026",
        "Sylvie DUMONT — Gestionnaire Paie",
        "Catherine LEFEVRE — Directrice RH",
        [
            ("h2", "1. Plafonds Remboursement Repas"),
            ("p", "• Déplacement en Province : Plafond maximum de 22,00 € TTC par repas."),
            ("p", "• Déplacement en Île-de-France et Métropoles (> 500k hab.) : Plafond maximum de 28,50 € TTC par repas."),
            ("p", "• Déplacement International (Export Allemagne/Benelux suivi par Hugo VIDAL) : Plafond de 45,00 € TTC."),
            ("p", "Facture originale avec TVA détaillée obligatoire (les tickets CB sans détail sont rejetés par Ahmed ZOUARI)."),
            ("h2", "2. Indemnités Kilométriques (IK 2026)"),
            ("table", [
                ["Puissance Fiscale", "Jusqu'à 5 000 km", "De 5 001 à 20 000 km", "Au-delà de 20 000 km"],
                ["3 et 4 CV", "0.415 € / km", "(Distance × 0.270) + 725 €", "0.330 € / km"],
                ["5 CV", "0.451 € / km", "(Distance × 0.315) + 680 €", "0.370 € / km"],
                ["6 CV et plus", "0.485 € / km", "(Distance × 0.340) + 725 €", "0.405 € / km"]
            ]),
            ("p", "Majoration Véhicule Électrique : +20 % sur le barème calculé."),
            ("h2", "3. Astreintes & Télétravail"),
            ("p", "• Prime forfaitaire astreinte nuit semaine : 45,00 € brut (SOP_RH_003)."),
            ("p", "• Prime forfaitaire astreinte week-end complet : 185,00 € brut."),
            ("p", "• Indemnité journalière télétravail : 2,60 € net / jour (plafond mensuel de 57,20 € net).")
        ]
    )

    create_docx(
        f"{BASE_DIR}/02_Guides_Calculs_et_Matrices/AIDE_CALCUL_05_TRS_Taux_Rendement_Synthetique_Lignes.docx",
        "Calcul et Suivi du Taux de Rendement Synthétique (TRS) Lignes",
        "AIDE-CALCUL-05", "1.0", "15 janvier 2026",
        "Camille OLLIVIER — Ingénieure Amélioration Continue",
        "Marc-Antoine PELLETIER — Directeur Industriel",
        [
            ("1. Formule du TRS Usine", [
                "Formule standard : TRS = Taux de Disponibilité (D) × Taux de Performance (P) × Taux de Qualité (Q)",
                "• Disponibilité (D) = [ Temps Fonctionnement Réel / Temps Ouverture Prévu ] × 100",
                "• Performance (P) = [ Quantité Produite Réelle / Quantité Théorique Nominale ] × 100",
                "• Qualité (Q) = [ Quantité Conforme Vendable / Quantité Totale Produite ] × 100"
            ]),
            ("2. Seuils et Alertes Usines", [
                [
                    ["Niveau d'Alerte", "Intervalle TRS", "Action Obligatoire", "Responsable"],
                    ["Objectif Cible", "TRS ≥ 82.0 %", "Fonctionnement nominal", "Chefs d'équipe (J. THIBAULT / A. FERRARI)"],
                    ["Vigilance", "78.0 % ≤ TRS < 82.0 %", "Analyse en réunion de production hebdo", "Samuel KONE / Chloé RENARD"],
                    ["Critique", "TRS < 78.0 %", "Déclenchement groupe 8D sous 48h", "Vincent LACROIX & Camille OLLIVIER"]
                ]
            ])
        ]
    )

    df_freinte = pd.DataFrame({
        "Code_Batch": ["BAT-2026-001", "BAT-2026-002", "BAT-2026-003", "BAT-2026-004", "BAT-2026-005"],
        "Produit": ["Pancakes Vanille x8", "Pancakes Myrtille x8", "Brioche Feuilletée 400g", "Brioche Tranchée 500g", "Sauce Vegan 250ml"],
        "Ligne_Fabrication": ["Ligne 1 (Amiens)", "Ligne 1 (Amiens)", "Ligne 2 (Amiens)", "Ligne 2 (Amiens)", "Ligne 3 (Avignon)"],
        "Masse_Pate_Crue_kg": [120.0, 125.0, 200.0, 250.0, 500.0],
        "Masse_Apres_Cuisson_kg": [105.6, 108.75, 178.0, 222.5, 495.0],
        "Perte_Masse_Cuisson_kg": [14.4, 16.25, 22.0, 27.5, 5.0],
        "Taux_Freinte_Reel_Pct": [12.0, 13.0, 11.0, 11.0, 1.0],
        "Seuil_Max_Tolerance_Pct": [12.5, 13.5, 11.5, 11.5, 1.2],
        "Statut_Conformite": ["CONFORME", "CONFORME", "CONFORME", "CONFORME", "CONFORME"]
    })
    create_excel(
        f"{BASE_DIR}/02_Guides_Calculs_et_Matrices/AIDE_CALCUL_02_Taux_de_Freinte_et_Pertes_Cuisson.xlsx",
        {"Suivi_Freinte_Cuisson": df_freinte}
    )

    df_cout_revient = pd.DataFrame({
        "Reference_SKU": ["SKU-PAN-01", "SKU-BRI-02", "SKU-SAU-03", "SKU-SAB-04"],
        "Designation": ["Pancakes Vanille x8 (280g)", "Brioche Feuilletée 400g", "Sauce Vegan 250ml", "Sablés Pur Beurre x12"],
        "Site_Fabrication": ["Amiens (Nord)", "Amiens (Nord)", "Avignon (Sud)", "Amiens (Nord)"],
        "Cout_Matieres_Premieres_EUR": [0.42, 0.85, 0.38, 0.52],
        "Cout_Emballage_rPET_EUR": [0.15, 0.18, 0.22, 0.12],
        "Main_Oeuvre_Directe_MOD_EUR": [0.25, 0.35, 0.15, 0.20],
        "Frais_Generaux_Usine_EUR": [0.12, 0.16, 0.09, 0.10],
        "Cout_Revient_Complet_EUR": [0.94, 1.54, 0.84, 0.94],
        "Taux_Marge_Brute_Cible_Pct": [35.0, 32.0, 40.0, 38.0],
        "Prix_Vente_Calcule_HT_EUR": [1.45, 2.26, 1.40, 1.52],
        "Tarif_Negocie_Client_Alpha_EUR": [1.38, 2.18, 1.35, 1.46]
    })
    create_excel(
        f"{BASE_DIR}/02_Guides_Calculs_et_Matrices/GUIDE_CALCUL_03_Cout_de_Revient_et_Marge_Brute_B2B.xlsx",
        {"Calcul_Cout_Revient": df_cout_revient}
    )

    # ==========================================
    # 03_Recettes_et_Formulation
    # ==========================================
    create_pdf(
        f"{BASE_DIR}/03_Recettes_et_Formulation/REC_TECH_Brioche_Feuilletee_Industrielle_v2.1.pdf",
        "Brioche Feuilletée 400g Pur Beurre Industrielle",
        "REC-TECH-BRI-02", "2.1", "8 janvier 2026",
        "Yann PICARD — Ingénieur R&D Boulangerie",
        "Karim OUALI — Resp. Formulation Boulangerie",
        [
            ("h2", "1. Caractéristiques Produit Fini"),
            ("p", "• Code Article SKU : SKU-BRI-02 | Conditionnement : Barquette rPET sous atmosphère modifiée (N2 70% / CO2 30%)."),
            ("p", "• Poids net : 400g | DLC : 28 jours | Allégation : Pur Beurre sans conservateurs."),
            ("h2", "2. Formule de Composition (Pétrin Standard 100 kg Farine)"),
            ("table", [
                ["Ingrédient", "Code MP", "Fournisseur Agréé", "% Boulanger", "Poids (kg)"],
                ["Farine Blé T45 Gruau", "MP-FAR-001", "Fournisseur A SAS (Secours: Fournisseur B)", "100.0 %", "100.00 kg"],
                ["Beurre Tourage 84%", "MP-BEU-004", "Fournisseur D Laiterie", "35.0 %", "35.00 kg"],
                ["Oeufs entiers past.", "MP-OEU-002", "Fournisseur C Ingrédients", "30.0 %", "30.00 kg"],
                ["Eau filtrée osmosée", "MP-EAU-000", "Réseau Interne Amiens", "28.0 %", "28.00 kg"],
                ["Sucre semoule fin", "MP-SUC-001", "Fournisseur E Sucres", "15.0 %", "15.00 kg"],
                ["Levure boulangère liquide", "MP-LEV-001", "Fournisseur F Levures", "4.5 %", "4.50 kg"],
                ["Sel fin de mer", "MP-SEL-001", "Fournisseur G Sel", "2.0 %", "2.00 kg"],
                ["Extrait Vanille Bourbon", "MP-VAN-001", "Fournisseur V Export", "0.8 %", "0.80 kg"]
            ]),
            ("h2", "3. Paramètres Ligne 2 (Amiens - Samuel KONE)"),
            ("p", "• Pétrissage : 4 min en vitesse 1, puis 8 min 30 s en vitesse 2. Température cible : 23.5°C ± 0.5°C."),
            ("p", "• Pousse : 75 min à 28°C, hygrométrie 80%."),
            ("p", "• Cuisson four tunnel continu : 18 min à 185°C. Injection de buée 45 secondes à l'enfournement.")
        ]
    )

    create_pdf(
        f"{BASE_DIR}/03_Recettes_et_Formulation/REC_TECH_Sauce_Emulsionnee_Vegan_v1.0.pdf",
        "Sauce Émulsionnée Froide Végétale 250ml",
        "REC-TECH-SAU-03", "1.0", "14 janvier 2026",
        "Manon DUBUISSON — Ingénieure R&D Sauces",
        "Elodie FONTAINE — Resp. Formulation Sauces",
        [
            ("h2", "1. Spécifications Produit Fini"),
            ("p", "• Code Article SKU : SKU-SAU-03 | Ligne de fabrication : Ligne 3 Usine Sud Avignon (Théo GARNIER)."),
            ("p", "• Formulation 100% Vegan garantie sans allergènes majeurs lait/œufs/gluten (MATRICE_Allergenes)."),
            ("h2", "2. Formule Batch 500 kg"),
            ("table", [
                ["Ingrédient", "Rôle Technologique", "Fournisseur", "Pourcentage", "Masse (kg)"],
                ["Huile Colza Désodorisée", "Phase lipidique", "Fournisseur H Huiles", "65.0 %", "325.0 kg"],
                ["Eau déminéralisée", "Phase aqueuse", "Réseau Avignon", "22.5 %", "112.5 kg"],
                ["Protéine Féverole Isolée", "Émulsifiant végétal", "Fournisseur C Ingrédients", "3.5 %", "17.5 kg"],
                ["Moutarde de Dijon standard", "Arôme & Stabilisation", "Fournisseur K", "4.0 %", "20.0 kg"],
                ["Vinaigre Alcool 10°", "Acidification", "Fournisseur L", "3.0 %", "15.0 kg"],
                ["Gomme Xanthane (E415)", "Stabilisant texture", "Fournisseur M", "0.3 %", "1.5 kg"],
                ["Sel fin et aromates", "Assaisonnement", "Fournisseur G", "1.7 %", "8.5 kg"]
            ])
        ]
    )

    df_allulose = pd.DataFrame({
        "Molecule_Substitut": ["Allulose Pur Cristallisé", "Érythritol Poudre", "Maltitol Sirop", "Extrait Stévia Reb M 95%", "Sucralose Haute Pureté"],
        "Pouvoir_Sucrant_vs_Saccharose": [0.70, 0.65, 0.90, 250.0, 600.0],
        "Valeur_Energetique_kcal_g": [0.2, 0.0, 2.4, 0.0, 0.0],
        "Indice_Glycemique": [0, 0, 35, 0, 0],
        "Caramelisation_Cuisson": ["Excellente (brunit comme saccharose)", "Nulle", "Moyenne", "Nulle", "Nulle"],
        "Statut_Reglementaire_UE": ["Dossier Novel Food EFSA en cours", "Autorisé", "Autorisé (>10% mention laxatif)", "Autorisé", "Autorisé"],
        "Fournisseur_Agréé": ["Fournisseur C Ingrédients", "Fournisseur C Ingrédients", "Fournisseur E Sucres", "Fournisseur C Ingrédients", "Fournisseur E Sucres"]
    })
    create_excel(
        f"{BASE_DIR}/03_Recettes_et_Formulation/MATRICE_Equivalence_Edulcorants_Allulose.xlsx",
        {"Equivalences_Edulcorants": df_allulose}
    )

    # ==========================================
    # 04_Packaging_et_DevPack
    # ==========================================
    create_pdf(
        f"{BASE_DIR}/04_Packaging_et_DevPack/SPEC_TECH_Operculage_Barquettes_rPET.pdf",
        "Spécifications Barquettes rPET et Film d'Operculage",
        "SPEC-TECH-EMB-01", "2.0", "16 janvier 2026",
        "Léa SIMON — Ingénieure Packaging",
        "Thomas WEBER — Resp. Packaging & DevPack",
        [
            ("h2", "1. Spécifications Techniques"),
            ("p", "• Barquette : rPET thermoformé issu à 80% de recyclage post-consommation, fourni par Fournisseur B Emballages (David ROUX)."),
            ("p", "• Épaisseur fond barquette : 450 micromètres (tolérance ± 20 µm)."),
            ("p", "• Film d'operculage : Complexe barrière multicouche PET/PE-EVOH pelable de 52 micromètres."),
            ("h2", "2. Réglages Machine Operculeuse Ligne 3"),
            ("p", "• Température des mâchoires : 165°C ± 3°C | Temps de scellage : 1,2 s sous pression de 4,5 bars."),
            ("p", "• Gaz de balayage : 70% Azote (N2) / 30% CO2. Taux d'O2 résiduel visé : < 0,5% (conforme TEST_Etancheite_et_Permeabilite_O2).")
        ]
    )

    create_docx(
        f"{BASE_DIR}/04_Packaging_et_DevPack/TEST_Etancheite_et_Permeabilite_O2.docx",
        "Rapport de Test d'Étanchéité et Perméabilité O2",
        "RAPPORT-TEST-DEV-2026-08", "1.0", "19 janvier 2026",
        "Léa SIMON — Ingénieure Packaging",
        "Nicolas BERTIN — Resp. Assurance Qualité",
        [
            ("1. Conditions de Contrôle Laboratoire", [
                "Échantillons : 100 barquettes rPET issues de l'essai pilote sur l'Usine Nord.",
                "Méthode : Test d'immersion sous cloche à vide (-400 mbar) et mesure coulométrique de l'OTR selon norme ASTM F1927."
            ]),
            ("2. Résultats & Conclusion", [
                "Taux de fuite constaté : 0% à -400 mbar pendant 60 secondes.",
                "Perméabilité OTR moyenne : 1,8 cm3/m2/24h (spécification max autorisée : 2,5 cm3/m2/24h).",
                "Conclusion : Emballage validé pour le conditionnement de la brioche feuilletée (REC_TECH_Brioche_Feuilletee_Industrielle_v2.1)."
            ])
        ]
    )

    create_pptx(
        f"{BASE_DIR}/04_Packaging_et_DevPack/REVUE_PROJET_Packaging_BioSource_2026.pptx",
        [
            ("Revue de Projet : Packaging Bio-Sourcé 2026", [
                "Équipe DevPack GUSTOO",
                "Thomas WEBER & Léa SIMON",
                "Site d'Amiens & Avignon"
            ]),
            ("Évaluation des Matériaux Alternatifs", [
                "Bagasse de canne à sucre : Rigidité validée mais perméabilité O2 incompatible avec DLC 28j.",
                "Cellulose moulée traitée : Validée pour biscuits secs, surcoût de +22% vs rPET.",
                "Barquette Carton / Liner PE pelable : Meilleur compromis testé avec Fournisseur B."
            ]),
            ("Jalons 2026", [
                "T2 2026 : Essai pilote Ligne 1 selon GUIDELINES_RD_019.",
                "T3 2026 : Audit de recyclabilité Citeo.",
                "T4 2026 : Déploiement sur la gamme Pancakes (Antoine JOUBERT)."
            ])
        ]
    )

    # ==========================================
    # 05_Fournisseurs_et_Achats
    # ==========================================
    df_fournisseurs = pd.DataFrame({
        "Code_Fournisseur": ["FOURN-A", "FOURN-B", "FOURN-C", "FOURN-D", "FOURN-E", "FOURN-F", "FOURN-H", "FOURN-V"],
        "Raison_Sociale": ["Fournisseur A SAS", "Fournisseur B Emballages", "Fournisseur C Ingrédients", "Fournisseur D Laiterie", "Fournisseur E Sucres", "Fournisseur F Levures", "Fournisseur H Huiles", "Fournisseur V Export"],
        "Acheteur_Referent": ["Claire VASSEUR", "David ROUX", "Claire VASSEUR", "Claire VASSEUR", "Claire VASSEUR", "Claire VASSEUR", "Claire VASSEUR", "Claire VASSEUR"],
        "Matiere_Principale": ["Farine Blé T45 Gruau", "Barquettes rPET 450µm", "Protéine Féverole & Allulose", "Beurre Tourage 84%", "Sucre Cristal Semoule", "Levure fraîche liquide", "Huile Colza Désodorisée", "Extrait Vanille Bourbon"],
        "Pays_Origine": ["France", "Belgique", "Pays-Bas", "France", "France", "France", "Allemagne", "Madagascar"],
        "Prix_Unitaire_HT_EUR": [0.85, 0.15, 4.20, 7.50, 1.10, 2.30, 1.65, 142.00],
        "Unite": ["kg", "unité", "kg", "kg", "kg", "kg", "kg", "Litre"],
        "Delai_Livraison_Jours": [5, 12, 8, 3, 4, 2, 7, 25],
        "Franco_Port_EUR": [1500, 3000, 2000, 1000, 1200, 500, 2500, 5000],
        "Certification": ["IFS Food v8", "BRCGS Packaging", "FSSC 22000", "IFS Food v8", "ISO 9001", "FSSC 22000", "IFS Food v8", "Ecocert Bio"],
        "Note_Audit_Performance": [94, 88, 91, 96, 85, 92, 90, 89]
    })
    create_excel(
        f"{BASE_DIR}/05_Fournisseurs_et_Achats/BASE_FOURNISSEURS_Matieres_Premieres_2026.xlsx",
        {"Matrice_Fournisseurs_2026": df_fournisseurs}
    )

    create_pdf(
        f"{BASE_DIR}/05_Fournisseurs_et_Achats/CONTRAT_CADRE_Fournisseur_Vanille_Madagascar.pdf",
        "Contrat-Cadre d'Approvisionnement : Extrait de Vanille Bourbon",
        "CONTRAT-ACH-2026-V", "1.0", "5 janvier 2026",
        "Claire VASSEUR — Acheteuse Matières Premières",
        "Philippe MOREL — Directeur Achats",
        [
            ("h2", "1. Objet et Engagement"),
            ("p", "Fourniture exclusive d'extrait liquide pur de vanille Bourbon de Madagascar certifié sans vanilline de synthèse pour les besoins des usines d'Amiens et Avignon."),
            ("h2", "2. Conditions Commerciales et Pénalités"),
            ("p", "• Volume annuel ferme : 4 500 Litres livrés en bidons PEHD de 25L."),
            ("p", "• Prix contractuel : 142,00 € HT le litre rendu DDP Amiens."),
            ("p", "• Pénalité de retard : 1,5% de la valeur de la commande par jour ouvré de retard de livraison.")
        ]
    )

    df_audit_huiles = pd.DataFrame({
        "Critere_Evalue": ["Hygiène des silos de stockage", "Traçabilité des lots", "Plan de contrôle 3-MCPD", "Gestion des alertes", "Ponctualité OTIF"],
        "Ponderation_Pct": [20, 25, 25, 15, 15],
        "Score_Sur_20": [18, 19, 17, 16, 18],
        "Statut": ["CONFORME", "EXCELLENT", "CONFORME", "ACCEPTABLE", "EXCELLENT"],
        "Auditeur": ["Mehdi SLIMANI", "Mehdi SLIMANI", "Nicolas BERTIN", "Mehdi SLIMANI", "Claire VASSEUR"]
    })
    create_excel(
        f"{BASE_DIR}/05_Fournisseurs_et_Achats/AUDIT_QUALITE_Fournisseur_Huiles_Vegetales.xlsx",
        {"Synthese_Audit_Huiles": df_audit_huiles}
    )

    create_pptx(
        f"{BASE_DIR}/05_Fournisseurs_et_Achats/PRESENTATION_Audit_Performance_Fournisseur_A.pptx",
        [
            ("Revue Annuelle : Fournisseur A SAS", [
                "Direction Achats & Assurance Qualité GUSTOO",
                "Philippe MOREL & Claire VASSEUR",
                "Période évaluée : Exercice 2025"
            ]),
            ("Bilan Opérationnel Farine T45", [
                "Volume livré en 2025 : 4 200 tonnes en citerne vrac sur Amiens.",
                "Taux de service OTIF : 98,2%.",
                "Note globale d'audit : 94/100 (Statut Fournisseur Stratégique A+)."
            ])
        ]
    )

    # ==========================================
    # 06_Clients_et_Comptes_Cles
    # ==========================================
    df_clients = pd.DataFrame({
        "ID_Client": ["CLI-ALPHA", "CLI-BETA", "CLI-GAMMA", "CLI-DELTA", "CLI-EPSILON"],
        "Raison_Sociale": ["Client Alpha Distribution", "Client Bêta Foodservice", "Client Gamma Restauration", "Client Delta Supermarchés", "Client Epsilon Export"],
        "KAM_Referent": ["Xavier GAUTHIER", "Camille ESSOMBA", "Camille ESSOMBA", "Xavier GAUTHIER", "Hugo VIDAL"],
        "CA_Annuel_2025_kEUR": [18500, 12400, 8900, 4200, 3100],
        "Volume_Livre_Tonnes": [9800, 6200, 4500, 1900, 1400],
        "Remise_Fin_Annee_Pct": [4.5, 3.0, 2.5, 1.5, 2.0],
        "Seuil_Volume_Remise_Tonnes": [8000, 5000, 4000, 1500, 1000],
        "Delai_Paiement_Jours": [60, 45, 30, 45, 60],
        "Penalite_Rupture": ["2% commande", "1% par jour retard", "500 € par palette", "1.5% commande", "Selon CMR"]
    })
    create_excel(
        f"{BASE_DIR}/06_Clients_et_Comptes_Cles/FICHIER_CLIENTS_B2B_et_Tarifs_Negocies.xlsx",
        {"Accords_Commerciaux_B2B": df_clients}
    )

    create_pdf(
        f"{BASE_DIR}/06_Clients_et_Comptes_Cles/CAHIER_CHARGES_Client_GrandDistribution_Carrefour.pdf",
        "Cahier des Charges MDD : Client Alpha Distribution",
        "CDC-CLI-ALPHA-2026", "1.0", "10 janvier 2026",
        "Xavier GAUTHIER — KAM Grande Distribution",
        "Isabelle ROBIN — Directrice Commerciale",
        [
            ("h2", "1. Exigences Ingrédients et Recettes MDD"),
            ("p", "• Interdiction stricte : Huile de palme, colorants azoïques et conservateurs de synthèse."),
            ("p", "• Oeufs de poules élevées en plein air (Code 1) d'origine France obligatoires."),
            ("p", "• Teneur maximale en sucre : 18g pour 100g de produit fini sur les brioches et pancakes."),
            ("h2", "2. Palettisation & Logistique"),
            ("p", "• Palettes Europe 800x1200mm houssées, hauteur max 1,60m. Étiquette GS1-128 conforme.")
        ]
    )

    create_pptx(
        f"{BASE_DIR}/06_Clients_et_Comptes_Cles/PRESENTATION_COMITE_Innovation_Client_Nestle.pptx",
        [
            ("Comité Innovation B2B : Client Bêta Foodservice", [
                "GUSTOO R&D & Direction Commerciale",
                "Isabelle ROBIN & Sophie LAMBERT",
                "Février 2026"
            ]),
            ("Projets en Co-Développement", [
                "Pancakes protéinés pour le petit-déjeuner hôtelier (15g protéines / portion).",
                "Sauces émulsionnées en poches souples 1kg pour distributeurs automatiques.",
                "Plan de lancement : Salons professionnels Q3 2026."
            ])
        ]
    )

    create_pptx(
        f"{BASE_DIR}/06_Clients_et_Comptes_Cles/PRESENTATION_Revue_Comite_Client_Alpha.pptx",
        [
            ("Revue Annuelle de Compte : Client Alpha Distribution", [
                "Compte Clé N°1 GUSTOO (18,5 M€ CA)",
                "Xavier GAUTHIER (KAM) & Marine COSTA (ADV)",
                "Bilan 2025 & Perspectives 2026"
            ]),
            ("Performance & Engagements", [
                "Taux de service logistique atteint en 2025 : 99,1% (Objectif contrat > 98,5%).",
                "Seuil de volume atteint (9 800 T) déclenchant la remise de 4,5%.",
                "Lancement de la référence Brioche Feuilletée 400g dans 350 hypermarchés en avril 2026."
            ])
        ]
    )

    # ==========================================
    # 07_Qualite_HACCP_et_Reglementaire
    # ==========================================
    df_allergenes = pd.DataFrame({
        "Ligne_Production": ["Ligne 1 (Pancakes - Amiens)", "Ligne 2 (Brioches - Amiens)", "Ligne 3 (Sauces - Avignon)"],
        "Responsable_Ligne": ["Samuel KONE", "Samuel KONE", "Théo GARNIER"],
        "Gluten": ["OUI (Farine T45)", "OUI (Farine T45)", "NON (Absence garantie certifiée)"],
        "Oeufs": ["OUI (Ovoproduits pasteurisés)", "OUI (Oeufs entiers)", "NON (Interdit - Ligne 100% Vegan)"],
        "Lait_Lactose": ["OUI (Lait entier & poudre)", "OUI (Beurre 84%)", "NON (Interdit - Ligne 100% Vegan)"],
        "Fruits_a_Coque": ["OUI (Poudre amandes cycle B)", "NON (Interdit strict Ligne 2)", "NON (Interdit strict Ligne 3)"],
        "Arachides": ["NON (Interdit site Usine)", "NON (Interdit site Usine)", "NON (Interdit site Usine)"],
        "Soja": ["Traces possibles (farine)", "NON (Absence totale)", "OUI (Isolats de soja utilisés)"],
        "Moutarde": ["NON (Absence)", "NON (Absence)", "OUI (Moutarde de Dijon standard)"],
        "Protocole_Nettoyage": ["NEP Niveau 2 (1h30)", "NEP Niveau 3 acide/soude (2h45)", "Stérilisation vapeur 121°C (45 min)"]
    })
    create_excel(
        f"{BASE_DIR}/07_Qualite_HACCP_et_Reglementaire/MATRICE_Allergenes_et_Lignes_Production_2026.xlsx",
        {"Matrice_Allergenes_Lignes": df_allergenes}
    )

    create_pdf(
        f"{BASE_DIR}/07_Qualite_HACCP_et_Reglementaire/PROCEDURES_Rappel_Lot_Non_Conforme.pdf",
        "Guide d'Audit et Traitement des Non-Conformités Internes",
        "PROC-QUAL-009", "1.1", "11 janvier 2026",
        "Aïcha TOURE — Resp. HACCP",
        "Hélène GIRAUD — Directrice QSE",
        [
            ("h2", "1. Classification des Anomalies Qualité"),
            ("p", "• Mineure : Défaut visuel packaging sans impact sanitaire (dérogation chef d'équipe)."),
            ("p", "• Majeure : Écart freinte cuisson hors tolérance ou DLC erronée (arrêt ligne immédiat)."),
            ("p", "• Critique : Dépassement CCP microbiologique ou présence allergène non étiqueté -> Déclenchement obligatoire de la SOP_QUAL_012 sous 2 heures.")
        ]
    )

    create_docx(
        f"{BASE_DIR}/07_Qualite_HACCP_et_Reglementaire/REGISTRE_Audits_Sanitaires_Usine_Nord.docx",
        "Registre Officiel des Contrôles Sanitaires DDPP et Vétérinaires",
        "REG-DDPP-2026", "1.0", "14 janvier 2026",
        "Mehdi SLIMANI — Auditeur Qualité Interne",
        "Nicolas BERTIN — Resp. Assurance Qualité",
        [
            ("1. Synthèse de l'Inspection DDPP du 14 Novembre 2025", [
                "Inspecteurs : DDPP Somme (Site d'Amiens).",
                "Conclusion : Établissement conforme (Niveau d'hygiène : Très satisfaisant).",
                "Action soldée : Remplacement des joints du sas de stockage levures réalisé le 12 janvier 2026 par Rachid BENALI."
            ]),
            ("2. Suivi Microbiologique Annuel (Listeria & Salmonella)", [
                "Bilan 2025 : 1 040 prélèvements de surface réalisés par Sarah NOEL (Usine Nord) et Bruno LEGRAND (Usine Sud). 0 non-conformité constatée."
            ])
        ]
    )

    create_docx(
        f"{BASE_DIR}/07_Qualite_HACCP_et_Reglementaire/RAPPORT_Audit_Interne_Usine_Nord.docx",
        "Rapport d'Audit Qualité Interne IFS Food v8",
        "RAPPORT-AUDIT-2026-01", "1.0", "20 janvier 2026",
        "Mehdi SLIMANI — Auditeur Qualité Interne",
        "Hélène GIRAUD — Directrice QSE",
        [
            ("1. Synthèse de l'Audit Amiens", [
                "Lignes auditées : Ligne 1 Pancakes et Ligne 2 Brioches en présence d'Odile FABRE.",
                "Score global atteint : 96,4% de conformité aux exigences IFS Food v8.",
                "Point d'amélioration : Rappeler la consigne de changement de blouse lors du passage du secteur préparation vers le sas de conditionnement."
            ])
        ]
    )

    # ==========================================
    # 08_Finance_et_Controle_de_Gestion
    # ==========================================
    df_budget = pd.DataFrame({
        "Poste_Depense": ["Matières Premières Farines & Ingrédients", "Emballages & Barquettes rPET", "Masse Salariale Production & Maintenance", "Énergie (Gaz fours & Électricité)", "Maintenance Industrielle & Pièces", "Prestations Hygiène & Déchets", "Amortissements CAPEX Lignes"],
        "Budget_Prevu_2026_EUR": [14200000, 3800000, 8500000, 2100000, 1200000, 450000, 1650000],
        "Realise_2025_EUR": [13400000, 3650000, 8100000, 2450000, 1150000, 420000, 1500000],
        "Ecart_Montant_EUR": [800000, 150000, 400000, -350000, 50000, 30000, 150000],
        "Variation_Pct": [6.0, 4.1, 4.9, -14.3, 4.3, 7.1, 10.0]
    })
    create_excel(
        f"{BASE_DIR}/08_Finance_et_Controle_de_Gestion/BUDGET_Previsionnel_Usine_Nord_2026.xlsx",
        {"Budget_Usine_Nord_2026": df_budget}
    )

    df_marges_q4 = pd.DataFrame({
        "Client": ["Client Alpha", "Client Bêta", "Client Gamma", "Client Delta", "Client Epsilon"],
        "KAM": ["Xavier GAUTHIER", "Camille ESSOMBA", "Camille ESSOMBA", "Xavier GAUTHIER", "Hugo VIDAL"],
        "Chiffre_Affaires_Q4_EUR": [5100000, 3200000, 2300000, 1100000, 850000],
        "Marge_Brute_Degagee_EUR": [1683000, 1120000, 874000, 429000, 323000],
        "Taux_Marge_Brute_Pct": [33.0, 35.0, 38.0, 39.0, 38.0],
        "Volume_Livre_Q4_Tonnes": [2800, 1750, 1150, 520, 410]
    })
    create_excel(
        f"{BASE_DIR}/08_Finance_et_Controle_de_Gestion/RAPPORT_Marges_par_Client_Q4_2025.xlsx",
        {"Marges_Clients_Q4": df_marges_q4}
    )

    create_pdf(
        f"{BASE_DIR}/08_Finance_et_Controle_de_Gestion/NOTE_Investissement_Ligne_Automatisee.pdf",
        "Note de Cadrage CAPEX 2026 : Operculage Robotisé Amiens",
        "CAPEX-2026-01", "1.0", "12 janvier 2026",
        "Emma ROLLAND — Contrôleuse de Gestion Usine",
        "Nathalie PERRIN — Directrice DAF",
        [
            ("h2", "1. Justification Économique et Technique"),
            ("p", "• Investissement matériel validé au Comex : 1 450 000 € HT sur la Ligne 3 du site d'Amiens."),
            ("p", "• Objectif : Augmentation de cadence à 75 barquettes/min et réduction des rebuts film sous les 0,4%."),
            ("h2", "2. Indicateurs Financiers"),
            ("p", "• Temps de retour sur investissement (Payback) : 2,8 ans | Taux de Rentabilité Interne (TRI) : 18,4%.")
        ]
    )

    # ==========================================
    # 09_RH_et_Vie_Entreprise
    # ==========================================
    create_pdf(
        f"{BASE_DIR}/09_RH_et_Vie_Entreprise/POL_RH_Notes_Frais_et_Deplacements_2026.pdf",
        "Politique RH 2026 : Frais Professionnels et Missions",
        "POL-RH-2026", "2.0", "2 janvier 2026",
        "Fatou DIALLO — Chargée RH",
        "Catherine LEFEVRE — Directrice RH",
        [
            ("h2", "1. Principes et Délais"),
            ("p", "Les notes de frais doivent être transmises sur le portail RH avant le 5 de chaque mois M+1 pour validation par le N+1 et paiement avec la paie du mois."),
            ("h2", "2. Déplacements en Train et Avion"),
            ("p", "• Train obligatoire pour tout trajet d'une durée inférieure à 3h30 (Classe 2 pour les trajets < 4h)."),
            ("p", "• Les déplacements en avion nécessitent l'accord préalable de la Direction Générale (Camille ROUSSEL).")
        ]
    )

    create_pdf(
        f"{BASE_DIR}/09_RH_et_Vie_Entreprise/GUIDE_Onboarding_Nouveaux_Collaborateurs.pdf",
        "Livret d'Accueil du Nouveau Collaborateur GUSTOO",
        "LIVRET-ACCUEIL-2026", "1.0", "2 janvier 2026",
        "Baptiste MARTY — Resp. Recrutement & Formation",
        "Catherine LEFEVRE — Directrice RH",
        [
            ("h2", "1. Organisation du Travail"),
            ("p", "• Siège administratif : Plages fixes obligatoires de 9h30 à 12h00 et de 14h00 à 16h30."),
            ("p", "• Usines d'Amiens et Avignon : Organisation en 3x8 (Matin 5h-13h, Après-midi 13h-21h, Nuit 21h-5h)."),
            ("h2", "2. Avantages Sociaux GUSTOO"),
            ("p", "• Prise en charge mutuelle d'entreprise à 60%."),
            ("p", "• Titres-restaurant d'une valeur faciale de 9,50 € (part patronale : 5,70 €).")
        ]
    )

    df_grille_rh = pd.DataFrame({
        "Niveau_Coefficient": ["N1 (Coef 140)", "N2 (Coef 160)", "N3 (Coef 190)", "N4 (Coef 240)", "N5 (Cadre Coef 300)", "N6 (Cadre Sup Coef 400)"],
        "Poste_Type": ["Opérateur Conditionnement", "Conducteur de Ligne", "Technicien Maintenance", "Chef d'Équipe Production", "Ingénieur R&D / Packaging", "Directeur de Département"],
        "Salaire_Brut_Min_EUR": [1820, 2050, 2350, 2800, 3500, 5200],
        "Salaire_Brut_Max_EUR": [2100, 2400, 2900, 3400, 4600, 7500],
        "Prime_Objectif_Annuelle_Pct": [0.0, 3.0, 5.0, 8.0, 12.0, 20.0],
        "Statut_Cadre": ["NON", "NON", "NON", "NON", "OUI", "OUI"]
    })
    create_excel(
        f"{BASE_DIR}/09_RH_et_Vie_Entreprise/GRILLE_Salariale_Confidentielle_2026.xlsx",
        {"Grille_Salaires_2026": df_grille_rh}
    )

    # ==========================================
    # 10_IT_et_Juridique
    # ==========================================
    create_pdf(
        f"{BASE_DIR}/10_IT_et_Juridique/CHARTE_Securite_Informatique_et_Usage_IA.pdf",
        "Charte Sécurité des SI et Usage des Outils d'IA",
        "CHARTE-IT-05", "1.0", "12 janvier 2026",
        "Guillaume PASCAL — Resp. Infrastructure & Cyber",
        "Romain AUBERT — Directeur DSI",
        [
            ("h2", "1. Règles de Sécurité des Mots de Passe"),
            ("p", "• Longueur minimale de 12 caractères avec 4 types de caractères différents. Renouvellement tous les 90 jours."),
            ("h2", "2. Confidentialité des Données et Usage de l'IA"),
            ("p", "• Interdiction stricte d'injecter des données sensibles (recettes, prix fournisseurs, données salariales) dans des IA génératives grand public externes."),
            ("p", "• Seul le système RAG interne GUSTOO sécurisé est homologué pour le traitement des documents de l'entreprise.")
        ]
    )

    create_docx(
        f"{BASE_DIR}/10_IT_et_Juridique/CONTRAT_NDA_Type_Partenaire.docx",
        "Accord de Confidentialité et Non-Divulgation (NDA Type)",
        "NDA-TYPE-2026", "1.0", "6 janvier 2026",
        "Benoît CAILLAUX — Juriste Contrats & Affaires",
        "Alexandra FERNANDEZ — Directrice Juridique",
        [
            ("1. Définition des Informations Confidentielles", [
                "Sont réputées strictement confidentielles l'ensemble des formules de recettes, spécifications packaging rPET, paramètres de ligne et conditions tarifaires communiquées par GUSTOO au Partenaire."
            ]),
            ("2. Durée de l'Engagement", [
                "Les obligations de confidentialité prévues au présent accord demeurent en vigueur pendant une durée de cinq (5) années à compter de la date de signature."
            ])
        ]
    )

    print(f"\n🎉 L'ensemble des 25 documents et tableaux réels ont été créés avec succès dans '{BASE_DIR}/' !")

if __name__ == "__main__":
    generer_documents()