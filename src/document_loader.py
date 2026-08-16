import os
import pypdf
import pandas as pd
from docx import Document
from pptx import Presentation


class ProcesseurDocuments:
    """Processeur polyvalent pour l'extraction et le découpage de documents d'entreprise.

    Prend en charge les formats : .pdf, .docx, .xlsx, .pptx.

    Attributes:
        taille_chunk (int): Taille cible d'un extrait de texte (en caractères).
        recouvrement (int): Nombre de caractères partagés entre deux chunks consécutifs.
    """

    def __init__(self, taille_chunk: int = 1000, recouvrement: int = 150):
        """Initialise le processeur de documents.

        Args:
            taille_chunk (int, optional): Taille cible en caractères d'un chunk. Défaut: 1000.
            recouvrement (int, optional): Chevauchement entre chunks consécutifs. Défaut: 150.
        """
        self.taille_chunk = taille_chunk
        self.recouvrement = recouvrement

    # ==========================================
    # 1. EXTRACTEURS PAR TYPE DE FORMAT
    # ==========================================

    def _extraire_pdf(self, chemin_fichier: str) -> list[dict]:
        """Extrait le texte page par page depuis un PDF."""
        extraits = []
        with open(chemin_fichier, "rb") as f:
            lecteur = pypdf.PdfReader(f)
            for num_page, page in enumerate(lecteur.pages, start=1):
                texte = page.extract_text()
                if texte and texte.strip():
                    extraits.append({
                        "texte": texte.strip(),
                        "page": num_page
                    })
        return extraits

    def _extraire_docx(self, chemin_fichier: str) -> list[dict]:
        """Extrait les paragraphes et tableaux d'un document Word (.docx)."""
        doc = Document(chemin_fichier)
        lignes = []

        # Paragraphes
        for para in doc.paragraphs:
            if para.text.strip():
                lignes.append(para.text.strip())

        # Tableaux Word
        for table in doc.tables:
            lignes.append("\n--- Tableau ---")
            for row in table.rows:
                cellules = [cell.text.strip() for cell in row.cells]
                lignes.append(" | ".join(cellules))

        texte_global = "\n".join(lignes)
        return [{"texte": texte_global, "page": 1}] if texte_global.strip() else []

    def _extraire_pptx(self, chemin_fichier: str) -> list[dict]:
        """Extrait le texte slide par slide d'une présentation PowerPoint (.pptx)."""
        prs = Presentation(chemin_fichier)
        extraits = []

        for num_slide, slide in enumerate(prs.slides, start=1):
            textes_slide = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            textes_slide.append(paragraph.text.strip())
            
            if textes_slide:
                extraits.append({
                    "texte": "\n".join(textes_slide),
                    "page": num_slide
                })
        return extraits

    def _extraire_xlsx(self, chemin_fichier: str) -> list[dict]:
        """Convertit les feuilles et lignes Excel (.xlsx) en texte sémantique structuré."""
        extraits = []
        fichier_excel = pd.ExcelFile(chemin_fichier)

        for sheet_name in fichier_excel.sheet_names:
            df = pd.read_excel(chemin_fichier, sheet_name=sheet_name)
            if df.empty:
                continue

            lignes_texte = [f"### Feuille / Tableau : {sheet_name}"]
            for index, row in df.iterrows():
                elements = []
                for col in df.columns:
                    val = row[col]
                    if pd.notna(val):
                        elements.append(f"{col}: {val}")
                if elements:
                    lignes_texte.append(f"- Entrée {index + 1} : " + " | ".join(elements))

            texte_feuille = "\n".join(lignes_texte)
            extraits.append({
                "texte": texte_feuille,
                "page": sheet_name  # Utilise le nom de l'onglet comme repère
            })
        return extraits

    # ==========================================
    # 2. CHARGEMENT UNIFIÉ & DECOUPAGE (CHUNKING)
    # ==========================================

    def charger_document(self, chemin_fichier: str) -> list[dict]:
        """Identifie l'extension du fichier et extrait son contenu textuel."""
        _, extension = os.path.splitext(chemin_fichier.lower())

        if extension == ".pdf":
            return self._extraire_pdf(chemin_fichier)
        elif extension == ".docx":
            return self._extraire_docx(chemin_fichier)
        elif extension == ".pptx":
            return self._extraire_pptx(chemin_fichier)
        elif extension in [".xlsx", ".xls"]:
            return self._extraire_xlsx(chemin_fichier)
        else:
            return []

    def decouper_texte(self, texte: str) -> list[str]:
        """Découpe un texte long en chunks chevauchants."""
        if len(texte) <= self.taille_chunk:
            return [texte]

        chunks = []
        debut = 0
        while debut < len(texte):
            fin = debut + self.taille_chunk
            chunk = texte[debut:fin]
            chunks.append(chunk.strip())
            debut += self.taille_chunk - self.recouvrement
        return chunks

    def traiter_dossier(self, dossier_racine: str) -> list[dict]:
        """Scanne récursivement un dossier et retourne tous les chunks prêts pour ChromaDB.

        Args:
            dossier_racine (str): Chemin du dossier contenant l'ensemble des sous-dossiers.

        Returns:
            list[dict]: Liste des chunks formatés avec texte et métadonnées complètes.
        """
        extensions_valides = {".pdf", ".docx", ".pptx", ".xlsx", ".xls"}
        chunks_finaux = []

        for racine, _, fichiers in os.walk(dossier_racine):
            for fichier in fichiers:
                _, ext = os.path.splitext(fichier.lower())
                if ext in extensions_valides:
                    chemin_complet = os.path.join(racine, fichier)
                    pages_extraites = self.charger_document(chemin_complet)

                    for item in pages_extraites:
                        morceaux = self.decouper_texte(item["texte"])
                        for idx_chunk, morceau in enumerate(morceaux, start=1):
                            if morceau:
                                chunks_finaux.append({
                                    "texte": morceau,
                                    "metadata": {
                                        "source": fichier,
                                        "chemin_relatif": os.path.relpath(chemin_complet, dossier_racine),
                                        "page": str(item["page"]),
                                        "chunk_id": idx_chunk
                                    }
                                })
        return chunks_finaux